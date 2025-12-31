import tkinter as tk
from tkcalendar import DateEntry
import os
from tkinter import filedialog, messagebox, font, ttk
from PIL import Image, ImageTk
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from math import pi
from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Inches  # For specifying dimensions in inches
from matplotlib.figure import Figure
from persiantools.jdatetime import JalaliDate
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from ttkthemes import ThemedTk
import pandas as pd
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
import csv
import hashlib
import re
import platform
import convert_numbers
import base64
from pptx.util import Pt
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.backends import default_backend
import json
from cryptography.fernet import Fernet


users = {}
files = {}

root = tk.Tk()
root.withdraw()
style = tk.ttk.Style()
persian_font = ("B Nazanin", 16)
style.configure('Style1.TEntry', font=persian_font, padding=5, foreground='black', background='lightgrey')
style.configure('TButton', font=persian_font)
document = Document()
params = ['', 0 , '',[0,0,0,0,0]]
is_plotted = []
canvas = None
entries = [0,0,0,'1',0,0,'1',0,0,'1',0,0,'1',0,0]
engname = []
id = None

latin_to_persian_chars = {
    'a': 'ش', 'b': 'ذ', 'c': 'ز', 'd': 'ی', 'e': 'ث',
    'f': 'ب', 'g': 'ل', 'h': 'ا', 'i': 'ه', 'j': 'ت',
    'k': 'ن', 'l': 'م', 'm': 'ئ', 'n': 'د', 'o': 'خ',
    'p': 'ح', 'q': 'ض', 'r': 'ق', 's': 'س', 't': 'ف',
    'u': 'ع', 'v': 'ر', 'w': 'ص', 'x': 'ط', 'y': 'غ',
    'z': 'ظ'
}
def translate_persian(event):
    widget = event.widget
    if event.char in latin_to_persian_chars:
        # Replace the character with its Persian equivalent
        widget.insert(tk.INSERT, latin_to_persian_chars[event.char])
        return "break"
    
persian_to_latin_chars = {
    'ش': 'a', 'ذ': 'b', 'ز': 'c', 'ی': 'd', 'ث': 'e',
    'ب': 'f', 'ل': 'g', 'ا': 'h', 'ه': 'i', 'ت': 'j',
    'ن': 'k', 'م': 'l', 'ئ': 'm', 'د': 'n', 'خ': 'o',
    'ح': 'p', 'ض': 'q', 'ق': 'r', 'س': 's', 'ف': 't',
    'ع': 'u', 'ر': 'v', 'ص': 'w', 'ط': 'x', 'غ': 'y',
    'ظ': 'z'
}

def translate_latin(event):
    widget = event.widget
    if event.char in persian_to_latin_chars:
        widget.insert(tk.INSERT, persian_to_latin_chars[event.char])
        return "break"

def on_key_press(event):
    widget = event.widget
    # Dictionary mapping keycode values to Arabic digits
    keycode_to_digit = {
        48: '0', 49: '1', 50: '2', 51: '3', 52: '4',
        53: '5', 54: '6', 55: '7', 56: '8', 57: '9'
    }
    
    # Check if the keycode is one of the number keys
    if event.keycode in keycode_to_digit:
        # Insert the corresponding Arabic digit
        widget.insert(tk.INSERT, keycode_to_digit[event.keycode])
        # Prevent the default behavior
        return "break"

def hide_file(file_path):
    if platform.system() == 'Windows':
        # Hide file on Windows
        os.system(f'attrib +h "{file_path}"')
        hidden_file_path = file_path
    else:
        # Hide file on Unix-based systems
        directory, filename = os.path.split(file_path)
        hidden_file_path = os.path.join(directory, '.' + filename)
        os.rename(file_path, hidden_file_path)
    
    return hidden_file_path

def append_text_to_cell(table, search_text, append_text):
    # Iterate over all rows and cells in the table
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                if search_text in paragraph.text:
                    # If the search text is found, append the additional text
                    run = paragraph.add_run()
                    run.text = append_text
                    run.font.size = Pt(13)
                    return True
    return False


def add_image_after_text(slide, search_text, image_path, img_width, img_height):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if search_text in paragraph.text:
                # Calculate the position to place the image below the text
                left = shape.left + (shape.width - img_width) / 2
                top = shape.top + shape.height  # Place image just below the text
                # Add the picture to the slide at the calculated position
                slide.shapes.add_picture(image_path, left, top, img_width, img_height)
                return True
    return False


def process_presentation():
    if len(is_plotted) == 0:
        messagebox.showerror('Error','ابتدا نمودار را رسم کنید')
        return
    pptx_path = os.getcwd() + '/data/modified_presentation.pptx'
    image_path = 'images/save.png'
    prs = Presentation(pptx_path)
    text_updates = [
        ("نام و نام خانوادگی کودک"   ,  ' ' + params[0]  + ' '),
        ("سن کودک",  ' '  + convert_numbers.english_to_persian(params[1].split('.')[0])+'.'+ convert_numbers.english_to_persian(params[1].split('.')[1])),
        ("تاریخ ارزیابی", ' '  + convert_numbers.english_to_persian(params[2].split('-')[0])+'.'+ convert_numbers.english_to_persian(params[2].split('-')[1]) +'.'+ convert_numbers.english_to_persian(params[2].split('-')[2])),
        ("شناسه ", ' ' + convert_numbers.english_to_persian(id))
    ]

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table: 
                for search_text, append_text in text_updates:
                    append_text_to_cell(shape.table, search_text, append_text)


    slide_index = 0
    slide = prs.slides[slide_index]
    print('params', params)
    # Iterate through the shapes in the slide
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            # If the shape is a chart, print details about it
            chart = shape.chart

            chart_data = CategoryChartData()
            chart_data.categories = [ 'psychomotor', 'cognitive','emotional', 'social']
            chart_data.add_series('Column1', (params[3][1], params[3][0], params[3][3], params[3][2]))  # Replace with your actual data
            chart.replace_data(chart_data)        

    slide_index = 1
    slide = prs.slides[slide_index]
    engname.append(en_name.get())
    text_updates = [
        ('name:'  , ' ' + engname[-1]),
        ("age:", ' ' +  params[1] ),
        ("date:", ' ' +   params[2]),
        ("ID:", ' ' +   id),


    ]
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table: 
                for search_text, append_text in text_updates:
                    append_text_to_cell(shape.table, search_text, append_text)

    # Iterate through the shapes in the slide
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            # If the shape is a chart, print details about it
            chart = shape.chart

            chart_data = CategoryChartData()
            chart_data.categories = [ 'psychomotor', 'cognitive', 'emotional', 'social']
            chart_data.add_series(chart.series[0].name, chart.series[0].values)  # Replace with your actual data
            chart_data.add_series(params[0],  (params[3][1], params[3][0], params[3][3], params[3][2])) # Replace with your actual data

        # Replace the data in the chart
            chart.replace_data(chart_data)
            #prs.save(os.getcwd() + '/data/modified_presentation.pptx')
        current_directory = os.getcwd()

        # Navigate two directories up
        two_levels_up = os.path.abspath(os.path.join(current_directory, '..', '..'))

        # Now, navigate into another folder from this location
        target_folder = "data\\pptxs"
        final_path = two_levels_up +'/'+ target_folder

    file_path = filedialog.asksaveasfilename(initialdir=final_path,defaultextension=".pptx",
                                             initialfile= engname[-1],
                                             filetypes=[("power point files", "*.pptx"), ("All files", "*.*")])
    if file_path:
        prs.save(file_path)
        open_pptx(file_path)


def new_record():
    global canvas, id, loaded_key
    for i in range(15):
        if i == 0:
            entry_objects[i][0].configure(state='normal')
            entry_objects[i][1].configure(state='normal')
            entry_objects[i][0].delete(0, tk.END)
            entry_objects[i][1].delete(0, tk.END)
        elif i == 3 or i == 6 or i == 9 or i == 12 :
                entry_objects[i].set(1) 
        elif i == 2:
            entry_objects[i].delete(0, tk.END)
            entry_objects[i].insert(0, datetime.now().strftime("%m/%d/%y"))
        else:
            entry_objects[i].delete(0, tk.END)
            
        if canvas:
            canvas.get_tk_widget().destroy()
            canvas = None
        id_label.configure(text=f" ID: {'  '}")
        id = None
    params = []
    loaded_key = None



def load_entries_from_csv():
    global id_label, id
    current_directory = os.getcwd()

    # Navigate two directories up
    two_levels_up = os.path.abspath(os.path.join(current_directory, '..', '..'))

    # Now, navigate into another folder from this location
    target_folder = "data\\csvs\\"+ username
    final_path = two_levels_up +'/'+ target_folder
    file_path = filedialog.askopenfilename(initialdir=final_path ,filetypes=[("CSV files", "*.csv"), ("All files", "*.*")])
    if file_path:
        with open(file_path, mode='r' , encoding="utf-8") as file:
            reader = csv.reader(file)
            next(reader)  # Skip the header row
            for i, row in enumerate(reader):
                index, value = int(row[0]), (row[1])
                if i == 0:
                    entry_objects[i][0].configure(state='normal')
                    entry_objects[i][1].configure(state='normal')
                    entry_objects[i][0].delete(0, tk.END)
                    entry_objects[i][1].delete(0, tk.END)

                    entry_objects[i][0].insert(0, value.split('%')[0])
                    entry_objects[i][1].insert(0, value.split('%')[1])
                    entry_objects[i][0].configure(state='readonly')
                    entry_objects[i][1].configure(state='readonly')

                    engname.append(value.split('%')[1])
                elif i == 3 or i == 6 or i == 9 or i == 12 :
                        entry_objects[i].set(value) 
                elif i == 15:
                    id = value
                    id_label.configure(text=f" ID: {str(id)}")
                elif i == 16:
                    pass
                else:
                    entry_objects[i].delete(0, tk.END)
                    entry_objects[i].insert(0, value)
        
        params = []
        submit()
    

def open_csv(event):
    global id_label, id
        # This function is triggered when an item in the listbox is selected
    w = event.widget
    index = int(w.curselection()[0])
    value = w.get(index)
    current_directory = os.getcwd()

    # Navigate two directories up
    two_levels_up = os.path.abspath(os.path.join(current_directory, '..', '..'))

    # Now, navigate into another folder from this location
    #final_path = two_levels_up +'/'+ target_folder
    file_path = two_levels_up+'/data/csvs/'+ str(value[0])+'/'+ str(value[1])+'_'+value[2]+'_'+value[3]+'.csv'
    if file_path:
        with open(file_path, mode='r' , encoding="utf-8") as file:
            reader = csv.reader(file)
            next(reader)  # Skip the header row
            for i, row in enumerate(reader):
                index, value = int(row[0]), (row[1])
                if i == 0:
                    entry_objects[i][0].configure(state='normal')
                    entry_objects[i][1].configure(state='normal')
                    entry_objects[i][0].delete(0, tk.END)
                    entry_objects[i][1].delete(0, tk.END)

                    entry_objects[i][0].insert(0, value.split('%')[0])
                    entry_objects[i][1].insert(0, value.split('%')[1])
                    entry_objects[i][0].configure(state='readonly')
                    entry_objects[i][1].configure(state='readonly')

                    engname.append(value.split('%')[1])
                elif i == 3 or i == 6 or i == 9 or i == 12 :
                        entry_objects[i].set(value) 
                elif i == 15:
                    id = value
                    id_label.configure(text=f" ID: {str(id)}")
                elif i == 16:
                    pass
                else:
                    entry_objects[i].delete(0, tk.END)
                    entry_objects[i].insert(0, value)
        
        params = []
        submit()

def submit():
    global id, cipher_suite
    try:
        # Collect values from number inputs
        values = []
        for i in range(15):
            if hasattr(entries[i], 'get') :
                if entries[i].get() == '':
                    messagebox.showerror('Error','تمام مقادیر باید وارد شود')
                    return
        for i, entry in enumerate(entries[3:15]):
            if i%3 != 0:
                values.append(float(entry.get()))
            else:
                values.append(float(entry))
        #values = [float(entry.get()) for entry in entries[3:15]]
        # Collect values from string inputs
        name = str(entries[0].get())
        age = float(entries[1].get())
        over_7_maxs = [21, 36, 18, 18]
        under_7_maxs = [15, 27, 15, 15]
        if age >= 7:
            for i in range(4):
                if float(entries[3+i*3+1].get()) > over_7_maxs[i] or float(entries[3+i*3+1].get()) < 0:
                    messagebox.showerror('Error','اعداد وارد شده در بازه صحیح قرار ندارد. ')
                    return
                if float(entries[3+i*3+2].get()) > 6 or float(entries[3+i*3+2].get()) < 0:
                    messagebox.showerror('Error','اعداد وارد شده در بازه صحیح قرار ندارد. ')
                    return
        if age < 7:
            for i in range(4):
                if float(entries[3+i*3+1].get()) > under_7_maxs[i] or float(entries[3+i*3+1].get()) < 0:
                    messagebox.showerror('Error','اعداد وارد شده در بازه صحیح قرار ندارد. ')
                    return
                if float(entries[3+i*3+2].get()) > 6 or float(entries[3+i*3+2].get()) < 0:
                    messagebox.showerror('Error','اعداد وارد شده در بازه صحیح قرار ندارد. ')
                    return
        states = [values[0], values[3], values[6], values[9]]
        date_format = "%m/%d/%y"
        jalali_date = (entries[2].get())
        jalali_date = JalaliDate.to_jalali(datetime.strptime(entries[2].get(), date_format).year, datetime.strptime(entries[2].get(), date_format).month, datetime.strptime(entries[2].get(), date_format).day)
        params[0] = str(name)
        params[1] = str(age)
        params[2] = str(jalali_date)
        if age >= 7:
            vals = [((((values[4]*6)/over_7_maxs[1])+(values[5]*states[1]))/2)*0.8,
                    ((((values[1]*6)/over_7_maxs[0])+(values[2]*states[0]))/2)*0.8,
                    ((((values[7]*6)/over_7_maxs[2])+(values[8]*states[2]))/2)*0.8,
                    ((((values[10]*6)/over_7_maxs[3])+(values[11]*states[3]))/2)*0.8]
        if age < 7:
            vals = [((((values[4]*6)/under_7_maxs[1])+(values[5]*states[1]))/2)*0.8,
                    ((((values[1]*6)/under_7_maxs[0])+(values[2]*states[0]))/2)*0.8,
                    ((((values[7]*6)/under_7_maxs[2])+(values[8]*states[2]))/2)*0.8,
                    ((((values[10]*6)/under_7_maxs[3])+(values[11]*states[3]))/2)*0.8]
        params[3][0] = vals[0]
        params[3][1] = vals[1]
        params[3][2] = vals[2]
        params[3][3] = vals[3]
        if not id:
            id = generate_short_custom_id(name)

        current_directory = os.getcwd()
        # Navigate two directories up
        two_levels_up = os.path.abspath(os.path.join(current_directory, '..', '..'))

        # Now, navigate into another folder from this location
        target_folder = "data/csvs/"+ username
        final_path = two_levels_up +'/'+ target_folder
        print(final_path)
        file_path = final_path +'/'+ en_name.get() + '_'+ str(jalali_date) + '_'+ id +'.csv'
        plot_radar_chart(vals, name)
        id_label.configure(text=f" ID: {str(id)}")
        with open(file_path, mode='w', newline='',encoding="utf-8" ) as file:
            writer = csv.writer(file)
            writer.writerow(["Index", "Value"])
            writer.writerow([0, name + '%' + en_name.get()])
            writer.writerow([1, age])
            writer.writerow([2, entries[2].get()])
            for i, entry in enumerate(values):
                writer.writerow([i+3, entry])
            writer.writerow([15, str(id)])
            writer.writerow([16, str(username)])



    except ValueError  as e: 
        print(str(e)) 
        messagebox.showerror('Error','تمام مقادیر باید صحیح وارد شود')



def plot_radar_chart(values, name):
    global canvas
    labels = [ 'Cognitive', 'Psychomotor', 'Social', 'Emotional']
    num_vars = len(labels)
    chart_frame = tk.Frame(root, width=700, height=350)
    chart_frame.place(x=90, y=420)
    angles = [n / float(num_vars) * 2 * pi for n in range(num_vars)]
    values += values[:1]
    angles += angles[:1]

    #plt.style.use('dark_background')

    fig, ax = plt.subplots(figsize=(6, 3.5), subplot_kw={'polar': True})
    ax.fill(angles, values, color='blue', alpha=0.25)
    ax.plot(angles, values, color='blue', marker='o')
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels)
   #et_title("result", weight='bold', size=10, position=(0.5, 1.1),
    #             horizontalalignment='center', verticalalignment='center')
    ax.set_ylim(0,6)
    ax.figure.savefig(os.getcwd() +'/images/save.png')
    # Update the canvas with the new figure
    if canvas:
        try:
            canvas.get_tk_widget().destroy()
        except:
            pass
    canvas = FigureCanvasTkAgg(fig, master=chart_frame)
    canvas.draw()
    canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    is_plotted.append(1)
    

def on_close():
    """ This function is called when the window is closed. """
    plt.close('all')
    root.destroy()  # Destroy the main window

    
# File path for storing user credentials
credentials_file = 'data/user_credentials.csv'

def generate_short_custom_id(name):
    now = datetime.now()  # Capture the current date and time
    file_path = os.getcwd()+ '/data/last_id.txt'
    last_id = 0
    try:
        with open(file_path, 'r') as file:
            content = file.read().strip()  # Read and strip any extra whitespace
            if content:
                last_id = int(content) + 1
                print(f"Last ID: {last_id}")
            else:
                print("The file is empty or does not contain a valid integer.")
    except FileNotFoundError:
        print("The file does not exist.")
    except ValueError:
        print("The file content is not a valid integer.")
    # Access each component of the date and time
    second = str(now.second).zfill(2)# Pad second with leading zero if needed
    
    # Concatenate all parts to form a unique ID
    unique_id = str(last_id) + second
    os.remove(file_path)
    with open(file_path, mode='w' ) as file:
        file.write(str(last_id))
        file.close()
    return unique_id


def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()


def save_user(username, name, password, chk):
    with open(credentials_file, 'a', newline='', encoding="utf-8" ) as file:
        writer = csv.writer(file)
        writer.writerow([username, name, hash_password(password), chk_state.get()])
        #credentials_file = hide_file(credentials_file)


def user_exists(username):
    if not os.path.exists(credentials_file):
        return False
    with open(credentials_file, 'r') as file:
        reader = csv.reader(file)
        for row in reader:
            if row[0] == username:
                return True
    return False


def validate_user(username, password):
    global name_list, superuser
    if not os.path.exists(credentials_file):
        return False
    with open(credentials_file, 'r', encoding="utf-8") as file:
        reader = csv.reader(file)
        for row in reader:
            if row[0] == username and row[2] == hash_password(password):
                name_list = row[1]
                superuser = row[3]
                print(superuser)

                return True
    return False


def is_valid_username(username):
    return re.match("^[a-zA-Z]{4,}$", username) is not None


def is_valid_password(password):
    return re.match("^(?=.*[A-Za-z])(?=.*\d)[A-Za-z\d]{8,}$", password) is not None


# Function to handle user signup
def signup():
    global usernames
    username = entry_signup_username.get()
    name = entry_signup_name.get()
    password = entry_signup_password.get()
    
    if user_exists(username):
        messagebox.showerror("Error", "این نام کاربری از قبل وجود  دارد")

    if not is_valid_username(username):
        messagebox.showerror("Error", "نام کاربری حداقل باید ۴ حرفی باشد و فقط از حروف انگلیسی استفاده  شود")
        return

    if not is_valid_password(password):
        messagebox.showerror("Error", "رمزعبور باید حداقل ۸ حرفی باشد و از حروف و اعداد انگلیسی استفاده شود")
        return

    else:
        save_user(username, name, password, chk_state)
        current_directory = os.getcwd()

        # Navigate two directories up
        try:
            two_levels_up = os.path.abspath(os.path.join(current_directory, '..', '..'))

            # Now, navigate into another folder from this location
            target_folder = "/data/csvs/"+ username
            final_path = two_levels_up +'/'+ target_folder
            os.makedirs(final_path)
        except Exception as e:
            messagebox.showerror('Error', str(e))

        messagebox.showinfo("Success", "حساب با موفقیت ساخته شد")
        signup_window.destroy()
        login_window.deiconify()


def login():
    global username, user_key
    username = entry_login_username.get()
    password = entry_login_password.get()
    
    if validate_user(username, password):
        login_window.withdraw()
        open_main_app_window()
    else:
        messagebox.showerror("Login Error", "نام کاربری و یا رمز عبور اشتباه  است")



def open_signup_window():
    global signup_window, entry_signup_username, entry_signup_password, entry_signup_name, chk, chk_state 
    signup_window = tk.Toplevel(root)
    signup_window.title("Signup")
    signup_window.geometry("800x800")
    signup_window.geometry("+{}+{}".format(root.winfo_x(), root.winfo_y()))  # Place at the same position as root


    image_path = os.getcwd()+"/images/cli.png"
    # Setting the icon for the window using iconphoto (cross-platform)
    icon_image = tk.PhotoImage(file=image_path)
    root.iconphoto(True, icon_image)  # The 'True' parameter makes this icon used for all windows if more than one is created
    original_image = Image.open(image_path)
    resized_image = original_image.resize((80, 80), Image.Resampling.LANCZOS)
    photo = ImageTk.PhotoImage(resized_image)

    label_image = tk.Label(signup_window, image=photo)
    label_image.image = photo
    label_image.pack(side=tk.TOP, padx=0, pady=20)

    label_signup_username = tk.Label(signup_window, font=("Helvetica", 15), text="نام کاربری")
    label_signup_username.pack(pady=5)
    entry_signup_username = tk.Entry(signup_window)
    entry_signup_username.pack(pady=5)
    entry_signup_username.bind("<KeyPress>", translate_latin)


    label_signup_name = tk.Label(signup_window, font=("Helvetica", 15), text="نام و نام خانوادگی")
    label_signup_name.pack(pady=5)
    entry_signup_name = tk.Entry(signup_window)
    entry_signup_name.pack(pady=5)
    entry_signup_name.bind("<KeyPress>", translate_persian)


    label_signup_password = tk.Label(signup_window, font=("Helvetica", 15), text="رمز عبور")
    label_signup_password.pack(pady=5)
    entry_signup_password = tk.Entry(signup_window, show="*")
    entry_signup_password.pack(pady=5)

    chk_state = tk.BooleanVar()
    chk_state.set(False)  # Set check state

    # Create Checkbutton
    chk = tk.Checkbutton(signup_window, text='ادمین کل', var=chk_state)
    chk.pack()

    button_signup = tk.Button(signup_window, text="ایجاد حساب", command=signup)
    button_signup.pack(pady=10)
    button_to_login = tk.Button(signup_window, text="ورود", command=lambda: switch_window(signup_window, open_login_window))
    button_to_login.pack(pady=5)
    signup_window.protocol("WM_DELETE_WINDOW", on_close)



def open_login_window():
    global login_window, entry_login_username, entry_login_password, root 

    login_window = tk.Toplevel(root)
    login_window.title("Login")
    login_window.geometry("800x800")
    login_window.geometry("+{}+{}".format(root.winfo_x(), root.winfo_y()))  # Place at the same position as root

    image_path = os.getcwd()+"/images/cli.png"
    # Setting the icon for the window using iconphoto (cross-platform)
    icon_image = tk.PhotoImage(file=image_path)
    root.iconphoto(True, icon_image)  # The 'True' parameter makes this icon used for all windows if more than one is created
    
    original_image = Image.open(image_path)
    resized_image = original_image.resize((80, 80), Image.Resampling.LANCZOS)
    photo = ImageTk.PhotoImage(resized_image)

    label_image = tk.Label(login_window, image=photo)
    label_image.image = photo
    label_image.pack(side=tk.TOP, padx=0, pady=20)

    label_login_username = tk.Label(login_window, font=("Helvetica", 15), text="نام کاربری")
    label_login_username.pack(pady=5)
    entry_login_username = tk.Entry(login_window)
    entry_login_username.pack(pady=5)


    label_login_password = tk.Label(login_window, font=("Helvetica", 15), text="رمزعبور")
    label_login_password.pack(pady=5)
    entry_login_password = tk.Entry(login_window, show="*")
    entry_login_password.pack(pady=5)

    button_login = tk.Button(login_window, text="ورود", command=login)
    button_login.pack(pady=10)
    button_to_signup = tk.Button(login_window, text="ایجاد حساب", command=lambda: switch_window(login_window, open_signup_window))
    button_to_signup.pack(pady=5)
    login_window.protocol("WM_DELETE_WINDOW", on_close)



def switch_window(current_window, target_function):
    current_window.withdraw()
    target_function()


def logout():
    loaded_key = None
    global root
    root.destroy()
    root = tk.Tk()
    root.withdraw() 
    open_login_window() 


def open_pptx(file_path):
    if platform.system() == "Windows":
        os.startfile(file_path)
    elif platform.system() == "Darwin":  # macOS
        os.system(f'open "{file_path}"')
    else:  # Linux and other Unix-like systems
        os.system(f'xdg-open "{file_path}"')


def get_deepest_folder_name(root_dir):
    deepest_folder = root_dir
    max_depth = 0

    for dirpath, dirnames, filenames in os.walk(root_dir):
        # Calculate the depth by counting the number of os.sep in the path
        depth = dirpath.count(os.sep)
        if depth > max_depth:
            max_depth = depth
            deepest_folder = dirpath

    return os.path.basename(deepest_folder)


def perform_search():
    # Base directory path
    current_directory = os.getcwd()

    # Navigate two directories up
    two_levels_up = os.path.abspath(os.path.join(current_directory, '..', '..'))

    # Now, navigate into another folder from this location
    target_folder = "data\\csvs\\"
    base_dir = two_levels_up +'/'+ target_folder
    
    # List to store all file paths
    data_files = []
    
    # List all entries in the directory and one level of subdirectories
    entries = os.listdir(base_dir)
    for entry in entries:
        entry_path = os.path.join(base_dir, entry)
        if os.path.isfile(entry_path):
            data_files.append(entry_path)  # Add file path to list
        elif os.path.isdir(entry_path):
            # List contents of the subdirectory
            sub_entries = os.listdir(entry_path)
            for sub_entry in sub_entries:
                sub_entry_path = os.path.join(entry_path, sub_entry)
                if os.path.isfile(sub_entry_path):
                    data_files.append(sub_entry_path)  # Add subdirectory file path to list

    # Assuming s_entry is a tkinter Entry widget for search input
    search_query = s_entry.get().lower()
    
    # Assuming listbox is a tkinter Listbox widget for displaying results
    listbox.delete(0, tk.END)  # Clear previous results

    # Loop through files, assuming filenames are formatted as expected with underscores
    for file_path in data_files:
        filename = os.path.basename(file_path)  # Get the filename from the full path
        if search_query in filename.split('_')[2].lower():
            # Insert matching item details extracted from filename into the Listbox
            item_parts = filename.split('_')
            if len(item_parts) > 2:
                print(str(file_path.split('\\'+ get_deepest_folder_name(file_path))[0]).split('\\')[-1])
                print(get_deepest_folder_name(file_path))
                listbox.insert(tk.END, [str(file_path.split('\\'+ get_deepest_folder_name(file_path))[0]).split("\\")[-1] , item_parts[0], item_parts[1], item_parts[2].split('.')[0]])




# Function to open the main application window
def open_main_app_window():
    global main_app_frame, icon_photo, en_name, entry_objects, id_label, s_entry, listbox
    entry_objects = [[],None,None,None,None,None,None,None,None,None,None,None,None,None,None,]
    root.deiconify()
    root.title("استعداد یابی")
    root.geometry("800x800")


    image_path = os.getcwd()+"/images/cli.png"


    # Setting the icon for the window using iconphoto (cross-platform)
    icon_image = tk.PhotoImage(file=image_path)
    root.iconphoto(True, icon_image)  # The 'True' parameter makes this icon used for all windows if more than one is created
    '''
    original_image = Image.open(image_path)
    resized_image = original_image.resize((80, 80), Image.Resampling.LANCZOS)
    photo = ImageTk.PhotoImage(resized_image)

    label_image = tk.Label(root, image=photo)
    label_image.image = photo
    label_image.pack(side=tk.TOP, padx=0, pady=0)
    '''

    top_frame = tk.Frame(root)
    top_frame.pack(fill=tk.X, padx=10, pady=10)

    # Configure grid weights
        # Configure grid weights
    top_frame.columnconfigure(0, weight=1)
    top_frame.columnconfigure(1, weight=1)

    # Add a Logout button in the frame
    new_button = tk.Button(top_frame, text="رکورد جدید", command=new_record)
    new_button.grid(row=0, column=0, sticky='w', padx=50)

    logout_button = tk.Button(top_frame, text="خروج", command=logout)
    logout_button.grid(row=0, column=0, sticky='w', padx=(10, 40))

    # Display the username in the frame
    user_label = tk.Label(top_frame, text=f"کاربر: {name_list}", font=("Helvetica", 15), anchor='e')
    user_label.grid(row=0, column=1, sticky='e')

    id_label = tk.Label(top_frame, text=f" ID: {'  '}", font=("Helvetica", 15), anchor='e')
    id_label.grid(row=0, column=1, sticky='w')


    s_entry = tk.Entry(top_frame, width=15)
    s_entry.grid(row=0, column=0, sticky='e', padx=(10, 40))

    # Create a Button to trigger the search
    search_button = tk.Button(top_frame, text="جست و جو", command=perform_search)
    search_button.grid(row=0, column=0, sticky='e', padx=(10, 10))

    # Create a Listbox to display the search results
    listbox = tk.Listbox(root, width=40, height=3)
    listbox.pack(padx=10, pady=10)
    listbox.bind('<<ListboxSelect>>', open_csv)

    line = tk.Frame(root, height=2, bg="black", bd=0)
    line.place(x=20, y=50)
    line.pack(fill='x', padx=0, pady=0)
   
    # Create the main app frame
    main_app_frame = tk.Frame(root)
    main_app_frame.pack(fill=tk.BOTH, expand=True)

    labels = ['نام', 'سن', 'تاریخ ارزیابی', 'وضعیت حرکتی در همسالان', 'حرکتی-والدین', 'حرکتی-مربی', 'وضعیت شناختی در همسالان', 'شناختی-والدین ', 'شناختی-مربی',
            'وضعیت اجتماعی در همسالان', 'اجتماعی-والدین ', 'اجتماعی-مربی','وضعیت هیجانی در همسالان', 'هیجانی-والدین', 'مربی-هیجانی']

    for i, label in enumerate(labels):
        if i % 3 == 2:
            frame = tk.Frame(main_app_frame, height=20, width=20)
            frame.place(x=40, y= 30 + int(i/3)*40)
            
            if i == 2:
                date_entry = DateEntry(frame, width=7, background='darkblue', foreground='white', borderwidth=2)
                date_entry.pack(side=tk.LEFT, fill=tk.X)
                entries[i] = date_entry
                entry_objects[i] = (date_entry)
            else:
                entry = ttk.Entry(frame, width=8)
                entry.pack(side=tk.LEFT, fill=tk.X)
                entry.bind("<KeyPress>", on_key_press)
                entries[i] = entry
                entry_objects[i] = (entry)
            tk.Label(frame, font=("B Nazanin", 14), text=label).pack(side=tk.RIGHT, padx=20)
        elif i % 3 == 1:
            frame1 = tk.Frame(main_app_frame, height=20, width=20)
            frame1.place(x=240, y=30 + int(i/3)*40)

            entry = ttk.Entry(frame1, width=8)
            entry.pack(side=tk.LEFT, fill=tk.X)

            entry.bind("<KeyPress>", on_key_press)
            entries[i] = entry
            entry_objects[i] = entry

            tk.Label(frame1, font=("B Nazanin", 14), text=label).pack(side=tk.RIGHT, padx=30)
        elif i % 3 == 0:
            frame3 = tk.Frame(main_app_frame, height=20, width=20)
            frame3.place(x=465, y=30 + int(i/3)*40)

            if i != 0 :
                options = ['0.5', '1', '1.5']

                combobox = ttk.Combobox(frame3, values=options, width=3)
                combobox.pack(pady=20)
                combobox.set(options[1])
                combobox.place(x=10)
                def on_selection(event, combobox=combobox, i=i):
                    v = combobox.get()
                    entries[i] = v
                entry_objects[i] = (combobox)
                combobox.bind('<<ComboboxSelected>>', lambda event, c=combobox, idx=i: on_selection(event, c, idx))
                tk.Label(frame3, font=("B Nazanin", 14), text=label).pack(side=tk.RIGHT, padx=90)
            else:
                w = 9 if i == 0 else 6
                entry = ttk.Entry(frame3, width=w)
                entry.pack(side=tk.LEFT, fill=tk.X)
                entry.bind("<KeyPress>", translate_latin)
                tk.Label(frame3, font=("B Nazanin", 14), text=label+'(فارسی/انگلیسی)').pack(side=tk.RIGHT, padx=10)
                entry1 = ttk.Entry(frame3, width=w)
                entry1.pack(side=tk.LEFT, fill=tk.X)
                entry1.bind("<KeyPress>", translate_persian)

                entries[i] = entry1 
                entry_objects[i] = ([entry1, entry])
                en_name = entry
            

    # Buttons for actions
    style = ttk.Style()
    style.theme_use('clam')
    frame = tk.Frame(main_app_frame, height=600, width=400)
    frame.place(x=260, y=250)

    ttk.Button(frame, text='نمودار', command=submit).pack(side=tk.LEFT, padx=3, pady=2)
    ttk.Button(frame, text='فایل نتیجه', command=process_presentation).pack(side=tk.LEFT, padx=3, pady=2)
    ttk.Button(frame, text='بازکردن فایل', command=lambda: load_entries_from_csv()).pack(side=tk.LEFT, padx=3, pady=2)
    root.protocol("WM_DELETE_WINDOW", on_close)



if __name__ == "__main__":
    #root = tk.Tk()
    #root.withdraw()
    open_login_window()

    root.mainloop()

